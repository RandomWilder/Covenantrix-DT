"""
Document Processor
Handles document text extraction and preprocessing
"""
import logging
from typing import Optional, Dict, Any
from pathlib import Path
import mimetypes
import io

from core.exceptions import ProcessingError

logger = logging.getLogger(__name__)

# Try to import OCR service
try:
    from domain.integrations.ocr import OCRService
    OCR_SERVICE_AVAILABLE = True
except ImportError:
    OCR_SERVICE_AVAILABLE = False
    logger.warning("OCR service not available")

# Try to import optional dependencies
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyPDF2 not available - PDF support disabled")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX support disabled")

try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("PIL/pytesseract not available - OCR support disabled")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("openpyxl not available - Excel support disabled")

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PowerPoint support disabled")


class DocumentProcessor:
    """
    Document processing service for text extraction
    """
    
    # Supported formats
    TEXT_FORMATS = {'.txt', '.md', '.csv', '.json', '.xml', '.rtf'}
    PDF_FORMATS = {'.pdf'}
    DOCX_FORMATS = {'.docx', '.doc'}
    IMAGE_FORMATS = {'.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'}
    EXCEL_FORMATS = {'.xlsx', '.xls'}
    PPTX_FORMATS = {'.pptx', '.ppt'}
    
    SUPPORTED_FORMATS = TEXT_FORMATS | PDF_FORMATS | DOCX_FORMATS | IMAGE_FORMATS | EXCEL_FORMATS | PPTX_FORMATS
    
    def __init__(self, ocr_service: Optional[OCRService] = None):
        """
        Initialize document processor
        
        Args:
            ocr_service: Optional OCR service for fallback processing
        """
        self.logger = logging.getLogger(__name__)
        self.ocr_service = ocr_service
        self.ocr_used = False  # Track if OCR was used in last extraction
        self.logger.info("DocumentProcessor initialized")
    
    async def extract_text(
        self,
        file_content: bytes,
        filename: str,
        mime_type: Optional[str] = None
    ) -> str:
        """
        Extract text from document with OCR fallback
        
        Args:
            file_content: Raw file bytes
            filename: Original filename
            mime_type: MIME type (optional)
            
        Returns:
            Extracted text
        """
        # Reset OCR tracking for this extraction
        self.ocr_used = False
        
        file_ext = Path(filename).suffix.lower()
        
        # Detect MIME type if not provided
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)
        
        try:
            # Try standard extraction first
            extracted_text = await self._extract_text_standard(file_content, filename, file_ext, mime_type)
            
            # Validate content quality
            if self._should_use_ocr_fallback(extracted_text, filename, file_ext):
                self.logger.info(f"Content quality low, attempting OCR fallback for {filename}")
                try:
                    ocr_result = await self._extract_text_with_ocr(file_content, filename, mime_type)
                    if ocr_result and self.validate_content(ocr_result):
                        self.logger.info(f"OCR fallback successful for {filename}")
                        self.ocr_used = True
                        return ocr_result
                    else:
                        self.logger.warning(f"OCR fallback failed for {filename}, using standard extraction")
                except Exception as ocr_error:
                    self.logger.warning(f"OCR fallback failed for {filename}: {str(ocr_error)}")
            
            return extracted_text
            
        except Exception as e:
            # If standard extraction fails, try OCR as last resort
            if self._can_use_ocr(filename, file_ext, mime_type):
                self.logger.info(f"Standard extraction failed, attempting OCR for {filename}")
                try:
                    ocr_result = await self._extract_text_with_ocr(file_content, filename, mime_type)
                    if ocr_result:
                        self.logger.info(f"OCR rescue successful for {filename}")
                        self.ocr_used = True
                        return ocr_result
                except Exception as ocr_error:
                    self.logger.error(f"OCR rescue failed for {filename}: {str(ocr_error)}")
            
            self.logger.error(f"Text extraction failed for {filename}: {str(e)}")
            raise ProcessingError(
                f"Failed to extract text: {str(e)}",
                details={"filename": filename}
            )
    
    def _extract_text(self, content: bytes) -> str:
        """Extract text from plain text file"""
        # Try multiple encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                text = content.decode(encoding)
                self.logger.debug(f"Successfully decoded with {encoding}")
                return self._clean_text(text)
            except UnicodeDecodeError:
                continue
        
        raise ProcessingError("Failed to decode text with any known encoding")
    
    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    self.logger.warning(f"Failed to extract text from PDF page {page_num + 1}: {e}")
                    continue
            
            if not text_parts:
                raise ProcessingError("No text could be extracted from PDF")
            
            full_text = '\n\n'.join(text_parts)
            return self._clean_text(full_text)
            
        except Exception as e:
            raise ProcessingError(f"PDF extraction failed: {str(e)}")
    
    def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            docx_file = io.BytesIO(content)
            doc = DocxDocument(docx_file)
            
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
            
            if not text_parts:
                raise ProcessingError("No text could be extracted from DOCX")
            
            full_text = '\n'.join(text_parts)
            return self._clean_text(full_text)
            
        except Exception as e:
            raise ProcessingError(f"DOCX extraction failed: {str(e)}")
    
    def _extract_image_ocr(self, content: bytes) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(io.BytesIO(content))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            if not text.strip():
                raise ProcessingError("No text could be extracted from image")
            
            return self._clean_text(text)
            
        except Exception as e:
            raise ProcessingError(f"OCR extraction failed: {str(e)}")
    
    def _extract_excel(self, content: bytes) -> str:
        """Extract text from Excel file"""
        try:
            excel_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())
                    if row_text:
                        text_parts.append(' | '.join(row_text))
                
                text_parts.append("")  # Empty line between sheets
            
            if not text_parts:
                raise ProcessingError("No text could be extracted from Excel file")
            
            full_text = '\n'.join(text_parts)
            return self._clean_text(full_text)
            
        except Exception as e:
            raise ProcessingError(f"Excel extraction failed: {str(e)}")
    
    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint file"""
        try:
            pptx_file = io.BytesIO(content)
            presentation = pptx.Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_parts.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(' | '.join(row_text))
                
                text_parts.append("")  # Empty line between slides
            
            if not text_parts:
                raise ProcessingError("No text could be extracted from PowerPoint file")
            
            full_text = '\n'.join(text_parts)
            return self._clean_text(full_text)
            
        except Exception as e:
            raise ProcessingError(f"PowerPoint extraction failed: {str(e)}")
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        # Remove excessive whitespace
        text = '\n'.join(line.strip() for line in text.splitlines() if line.strip())
        
        # Normalize line endings
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        return text.strip()
    
    def validate_content(self, text: str) -> bool:
        """
        Validate extracted content quality
        
        Args:
            text: Extracted text
            
        Returns:
            True if content is valid
        """
        if not text or len(text) < 5:  # Lowered threshold
            return False
        
        # Check ratio of letters to total characters
        letters = sum(c.isalpha() for c in text)
        ratio = letters / len(text)
        
        # If less than 30% letters, likely gibberish or corrupted (lowered from 50%)
        return ratio >= 0.3
    
    async def _extract_text_standard(
        self,
        file_content: bytes,
        filename: str,
        file_ext: str,
        mime_type: Optional[str]
    ) -> str:
        """Extract text using standard methods"""
        # Text files
        if file_ext in self.TEXT_FORMATS:
            return self._extract_text(file_content)
        
        # PDF files
        elif file_ext in self.PDF_FORMATS:
            if not PDF_AVAILABLE:
                raise ProcessingError("PDF support not available - install PyPDF2")
            return self._extract_pdf(file_content)
        
        # DOCX files
        elif file_ext in self.DOCX_FORMATS:
            if not DOCX_AVAILABLE:
                raise ProcessingError("DOCX support not available - install python-docx")
            return self._extract_docx(file_content)
        
        # Image files (legacy OCR)
        elif file_ext in self.IMAGE_FORMATS:
            if not OCR_AVAILABLE:
                raise ProcessingError("OCR support not available - install PIL and pytesseract")
            return self._extract_image_ocr(file_content)
        
        # Excel files
        elif file_ext in self.EXCEL_FORMATS:
            if not EXCEL_AVAILABLE:
                raise ProcessingError("Excel support not available - install openpyxl")
            return self._extract_excel(file_content)
        
        # PowerPoint files
        elif file_ext in self.PPTX_FORMATS:
            if not PPTX_AVAILABLE:
                raise ProcessingError("PowerPoint support not available - install python-pptx")
            return self._extract_pptx(file_content)
        
        # Unsupported format
        raise ProcessingError(
            f"Unsupported file format: {file_ext}",
            details={"filename": filename, "mime_type": mime_type}
        )
    
    def _should_use_ocr_fallback(self, text: str, filename: str, file_ext: str) -> bool:
        """
        Determine if OCR fallback should be used
        
        Args:
            text: Extracted text
            filename: Original filename
            file_ext: File extension
            
        Returns:
            True if OCR fallback should be used
        """
        # Check if OCR service is available
        if not self.ocr_service or not OCR_SERVICE_AVAILABLE:
            return False
        
        # Check if OCR is enabled
        if not self.ocr_service.is_ocr_enabled():
            return False
        
        # Check if file format supports OCR
        if not self._can_use_ocr(filename, file_ext, None):
            return False
        
        # Check content quality
        if not self.validate_content(text):
            return True
        
        # Check minimum character count
        if len(text) < 5:  # Lowered threshold
            return True
        
        # Check for very low quality text (mostly non-letters)
        letters = sum(c.isalpha() for c in text)
        if len(text) > 0 and letters / len(text) < 0.3:
            return True
        
        return False
    
    def _can_use_ocr(self, filename: str, file_ext: str, mime_type: Optional[str]) -> bool:
        """
        Check if OCR can be used for this file
        
        Args:
            filename: Original filename
            file_ext: File extension
            mime_type: MIME type
            
        Returns:
            True if OCR can be used
        """
        if not self.ocr_service or not OCR_SERVICE_AVAILABLE:
            return False
        
        # Check if OCR is enabled
        if not self.ocr_service.is_ocr_enabled():
            return False
        
        # Check supported formats
        if file_ext in self.IMAGE_FORMATS:
            return self.ocr_service.validate_image_format(filename)
        elif file_ext in self.PDF_FORMATS:
            # PDFs can be processed with Google Vision API
            return True
        elif mime_type and mime_type.startswith('image/'):
            return True
        elif mime_type and mime_type == 'application/pdf':
            return True
        
        return False
    
    async def _extract_text_with_ocr(
        self,
        file_content: bytes,
        filename: str,
        mime_type: Optional[str]
    ) -> Optional[str]:
        """
        Extract text using OCR service
        
        Args:
            file_content: File content as bytes
            filename: Original filename
            mime_type: MIME type
            
        Returns:
            Extracted text or None if OCR fails
        """
        if not self.ocr_service:
            return None
        
        try:
            ocr_result = await self.ocr_service.extract_text(
                file_content=file_content,
                filename=filename,
                mime_type=mime_type
            )
            
            if ocr_result and ocr_result.text:
                self.logger.info(
                    f"OCR extracted {len(ocr_result.text)} characters "
                    f"with confidence {ocr_result.confidence:.2f}"
                )
                return ocr_result.text
            
            return None
            
        except Exception as e:
            self.logger.error(f"OCR extraction failed: {str(e)}")
            return None
    
    @staticmethod
    def is_supported_format(filename: str) -> bool:
        """Check if file format is supported"""
        ext = Path(filename).suffix.lower()
        return ext in DocumentProcessor.SUPPORTED_FORMATS